#!/usr/bin/env python3
"""
AI Presentation Generator - Python Script
Generates PowerPoint presentations from AI-generated outlines using python-pptx
"""

import json
import sys
import os
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import logging

# Configure logging to stderr only
logging.basicConfig(level=logging.INFO, stream=sys.stderr)
logger = logging.getLogger(__name__)

class PresentationGenerator:
    def __init__(self):
        self.templates = {
            'corporate_blue': {
                'primary_color': RGBColor(0, 51, 102),      # Dark blue
                'secondary_color': RGBColor(0, 102, 204),   # Light blue
                'accent_color': RGBColor(255, 255, 255),    # White
                'text_color': RGBColor(51, 51, 51),         # Dark gray
                'background_color': RGBColor(240, 248, 255), # Light blue background
                'title_color': RGBColor(0, 51, 102),        # Dark blue
                'subtitle_color': RGBColor(0, 102, 204)     # Light blue
            },
            'modern_white': {
                'primary_color': RGBColor(51, 51, 51),      # Dark gray
                'secondary_color': RGBColor(102, 102, 102), # Medium gray
                'accent_color': RGBColor(0, 0, 0),          # Black
                'text_color': RGBColor(51, 51, 51),         # Dark gray
                'background_color': RGBColor(255, 255, 255), # White
                'title_color': RGBColor(0, 0, 0),           # Black
                'subtitle_color': RGBColor(51, 51, 51)      # Dark gray
            },
            'creative_colorful': {
                'primary_color': RGBColor(255, 87, 34),     # Orange
                'secondary_color': RGBColor(33, 150, 243),  # Blue
                'accent_color': RGBColor(76, 175, 80),      # Green
                'text_color': RGBColor(51, 51, 51),         # Dark gray
                'background_color': RGBColor(255, 248, 240), # Light orange background
                'title_color': RGBColor(255, 87, 34),       # Orange
                'subtitle_color': RGBColor(33, 150, 243)    # Blue
            },
            'minimalist_gray': {
                'primary_color': RGBColor(97, 97, 97),      # Medium gray
                'secondary_color': RGBColor(158, 158, 158), # Light gray
                'accent_color': RGBColor(33, 33, 33),       # Dark gray
                'text_color': RGBColor(33, 33, 33),         # Dark gray
                'background_color': RGBColor(245, 245, 245), # Light gray background
                'title_color': RGBColor(33, 33, 33),        # Dark gray
                'subtitle_color': RGBColor(97, 97, 97)      # Medium gray
            },
            'academic_formal': {
                'primary_color': RGBColor(33, 33, 33),      # Dark gray/black
                'secondary_color': RGBColor(97, 97, 97),    # Medium gray
                'accent_color': RGBColor(255, 255, 255),    # White
                'text_color': RGBColor(33, 33, 33),         # Dark gray
                'background_color': RGBColor(255, 255, 255), # White
                'title_color': RGBColor(33, 33, 33),        # Dark gray
                'subtitle_color': RGBColor(97, 97, 97)      # Medium gray
            },
            'tech_modern': {
                'primary_color': RGBColor(0, 150, 136),     # Teal
                'secondary_color': RGBColor(76, 175, 80),   # Green
                'accent_color': RGBColor(255, 193, 7),      # Amber
                'text_color': RGBColor(33, 33, 33),         # Dark gray
                'background_color': RGBColor(240, 255, 250), # Light green background
                'title_color': RGBColor(0, 150, 136),       # Teal
                'subtitle_color': RGBColor(76, 175, 80)     # Green
            },
            'elegant_purple': {
                'primary_color': RGBColor(103, 58, 183),    # Purple
                'secondary_color': RGBColor(156, 39, 176),  # Deep purple
                'accent_color': RGBColor(255, 255, 255),    # White
                'text_color': RGBColor(51, 51, 51),         # Dark gray
                'background_color': RGBColor(248, 240, 255), # Light purple background
                'title_color': RGBColor(103, 58, 183),      # Purple
                'subtitle_color': RGBColor(156, 39, 176)    # Deep purple
            },
            'professional_green': {
                'primary_color': RGBColor(46, 125, 50),     # Dark green
                'secondary_color': RGBColor(76, 175, 80),   # Light green
                'accent_color': RGBColor(255, 255, 255),    # White
                'text_color': RGBColor(51, 51, 51),         # Dark gray
                'background_color': RGBColor(240, 255, 240), # Light green background
                'title_color': RGBColor(46, 125, 50),       # Dark green
                'subtitle_color': RGBColor(76, 175, 80)     # Light green
            }
        }
        
        self.font_styles = {
            'modern': {'font_name': 'Calibri', 'title_size': 44, 'content_size': 24},
            'classic': {'font_name': 'Times New Roman', 'title_size': 40, 'content_size': 22},
            'minimalist': {'font_name': 'Arial', 'title_size': 36, 'content_size': 20},
            'creative': {'font_name': 'Segoe UI', 'title_size': 42, 'content_size': 24}
        }

    def generate_presentation(self, data):
        """Main method to generate PowerPoint presentation"""
        try:
            logger.info("Starting presentation generation")
            logger.info(f"Received data: {data}")
            
            # Extract data
            outline = data.get('outline', {})
            template_name = data.get('template', 'corporate_blue')
            color_scheme = data.get('color_scheme', 'blue')
            font_style = data.get('font_style', 'modern')
            user_id = data.get('user_id')
            ai_result_id = data.get('ai_result_id')
            
            logger.info(f"Template name: {template_name}")
            logger.info(f"Available templates: {list(self.templates.keys())}")
            
            # Create new presentation
            prs = Presentation()
            
            # Apply template colors and fonts
            template_colors = self.templates.get(template_name, self.templates['corporate_blue'])
            font_config = self.font_styles.get(font_style, self.font_styles['modern'])
            
            # Generate slides
            slides_data = outline.get('slides', [])
            title = outline.get('title', 'Untitled Presentation')
        
            # Create title slide
            self.create_title_slide(prs, title, template_colors, font_config)
        
            # Create content slides with different layouts
            for i, slide_data in enumerate(slides_data):
                if slide_data.get('slide_type') == 'title':
                    continue  # Skip title slides as we already created one
                
                # Determine slide layout based on content and position
                content_items = slide_data.get('content', slide_data.get('subheaders', []))
                item_count = len(content_items) if isinstance(content_items, list) else 0
                
                # Use two-column layout for slides with many items
                if item_count > 4:
                    self.create_two_column_slide(
                        prs, 
                        slide_data, 
                        template_colors, 
                        font_config
                    )
                else:
                    # Use standard content slide
                    self.create_content_slide(
                        prs, 
                        slide_data, 
                        template_colors, 
                        font_config
                    )
            
            # Save presentation
            output_path = self.save_presentation(prs, user_id, ai_result_id)
            
            # Get actual file size
            file_size = 0
            if os.path.exists(output_path):
                file_size = os.path.getsize(output_path)
            
            return {
                'success': True,
                'file_path': output_path,
                'file_size': file_size,
                'download_url': f'/api/files/download/{os.path.basename(output_path)}',
                'slide_count': len(prs.slides)
            }
        
        except Exception as e:
            logger.error(f"Presentation generation failed: {str(e)}")
            return {
                'success': False,
                'error': str(e)
            }

    def add_decorative_elements(self, slide, colors):
        """Add decorative elements to title slide"""
        try:
            # Add corner accent shapes
            corner_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(8.5), Inches(6.5),
                Inches(1.5), Inches(1)
            )
            corner_fill = corner_shape.fill
            corner_fill.solid()
            corner_fill.fore_color.rgb = colors['accent_color']
            corner_shape.line.fill.background()
            
            # Add bottom accent line
            bottom_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(7),
                Inches(10), Inches(0.1)
            )
            bottom_fill = bottom_line.fill
            bottom_fill.solid()
            bottom_fill.fore_color.rgb = colors['secondary_color']
            bottom_line.line.fill.background()
            
        except Exception as e:
            logger.warning(f"Could not add decorative elements: {e}")

    def add_content_decorations(self, slide, colors):
        """Add decorative elements to content slides"""
        try:
            # Add corner accent
            corner_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(8.5), Inches(6.5),
                Inches(1.5), Inches(0.5)
            )
            corner_fill = corner_shape.fill
            corner_fill.solid()
            corner_fill.fore_color.rgb = colors['accent_color']
            corner_shape.line.fill.background()
            
            # Add side accent line
            side_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.2), Inches(1.5),
                Inches(0.1), Inches(5.5)
            )
            side_fill = side_line.fill
            side_fill.solid()
            side_fill.fore_color.rgb = colors['primary_color']
            side_line.line.fill.background()
            
        except Exception as e:
            logger.warning(f"Could not add content decorations: {e}")

    def create_two_column_slide(self, prs, slide_data, colors, fonts):
        """Create a two-column layout slide"""
        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background_color']
        
        # Add header bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0.5), 
            Inches(10), Inches(0.3)
        )
        header_fill = header_bar.fill
        header_fill.solid()
        header_fill.fore_color.rgb = colors['primary_color']
        header_bar.line.fill.background()
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get('header', 'Two Column Layout')
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.name = fonts['font_name']
        title_paragraph.font.size = Pt(fonts['title_size'])
        title_paragraph.font.color.rgb = colors['title_color']
        title_paragraph.font.bold = True
        
        # Create left column
        left_column = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(4.5), Inches(5.5)
        )
        left_fill = left_column.fill
        left_fill.solid()
        left_fill.fore_color.rgb = RGBColor(255, 255, 255)
        left_column.line.color.rgb = colors['secondary_color']
        left_column.line.width = Pt(2)
        
        # Create right column
        right_column = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(1.5),
            Inches(4.5), Inches(5.5)
        )
        right_fill = right_column.fill
        right_fill.solid()
        right_fill.fore_color.rgb = RGBColor(255, 255, 255)
        right_column.line.color.rgb = colors['secondary_color']
        right_column.line.width = Pt(2)
        
        # Add content to columns
        content_items = slide_data.get('content', slide_data.get('subheaders', []))
        if isinstance(content_items, list) and len(content_items) > 0:
            # Split content between columns
            mid_point = len(content_items) // 2
            
            # Left column content
            left_frame = left_column.text_frame
            left_frame.word_wrap = True
            left_frame.margin_left = Inches(0.2)
            left_frame.margin_right = Inches(0.2)
            left_frame.margin_top = Inches(0.2)
            left_frame.margin_bottom = Inches(0.2)
            
            for i, item in enumerate(content_items[:mid_point]):
                if item and item.strip():
                    if i == 0:
                        p = left_frame.paragraphs[0]
                    else:
                        p = left_frame.add_paragraph()
                    
                    text = item.strip()
                    if not text.startswith('•'):
                        text = f"• {text}"
                    
                    p.text = text
                    p.font.name = fonts['font_name']
                    p.font.size = Pt(fonts['content_size'])
                    p.font.color.rgb = colors['text_color']
                    p.level = 0
                    p.space_after = Pt(8)
            
            # Right column content
            right_frame = right_column.text_frame
            right_frame.word_wrap = True
            right_frame.margin_left = Inches(0.2)
            right_frame.margin_right = Inches(0.2)
            right_frame.margin_top = Inches(0.2)
            right_frame.margin_bottom = Inches(0.2)
            
            for i, item in enumerate(content_items[mid_point:]):
                if item and item.strip():
                    if i == 0:
                        p = right_frame.paragraphs[0]
                    else:
                        p = right_frame.add_paragraph()
                    
                    text = item.strip()
                    if not text.startswith('•'):
                        text = f"• {text}"
                    
                    p.text = text
                    p.font.name = fonts['font_name']
                    p.font.size = Pt(fonts['content_size'])
                    p.font.color.rgb = colors['text_color']
                    p.level = 0
                    p.space_after = Pt(8)
        
        # Add decorative elements
        self.add_content_decorations(slide, colors)

    def create_title_slide(self, prs, title, colors, fonts):
        """Create enhanced title slide with modern design"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background_color']
        
        # Add decorative header shape
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0), 
            Inches(10), Inches(1.5)
        )
        header_fill = header_shape.fill
        header_fill.solid()
        header_fill.fore_color.rgb = colors['primary_color']
        
        # Remove header shape outline
        header_shape.line.fill.background()
        
        # Set title with enhanced styling
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.name = fonts['font_name']
        title_paragraph.font.size = Pt(fonts['title_size'] + 8)  # Larger title
        title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text on colored background
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add subtitle placeholder with company info
        subtitle_shape = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), 
            Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.text = "Professional Presentation"
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.font.name = fonts['font_name']
        subtitle_paragraph.font.size = Pt(24)
        subtitle_paragraph.font.color.rgb = colors['secondary_color']
        subtitle_paragraph.font.italic = True
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add decorative elements
        self.add_decorative_elements(slide, colors)

    def create_content_slide(self, prs, slide_data, colors, fonts):
        """Create enhanced content slide with modern design"""
        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors['background_color']
        
        # Add header bar
        header_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0), Inches(0.5), 
            Inches(10), Inches(0.3)
        )
        header_fill = header_bar.fill
        header_fill.solid()
        header_fill.fore_color.rgb = colors['primary_color']
        header_bar.line.fill.background()
        
        # Set title with enhanced styling
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get('header', 'Untitled Slide')
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.name = fonts['font_name']
        title_paragraph.font.size = Pt(fonts['title_size'])
        title_paragraph.font.color.rgb = colors['title_color']
        title_paragraph.font.bold = True
        
        # Create content area with modern design
        content_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(9), Inches(5.5)
        )
        content_fill = content_shape.fill
        content_fill.solid()
        content_fill.fore_color.rgb = RGBColor(255, 255, 255)  # White content area
        content_shape.line.color.rgb = colors['secondary_color']
        content_shape.line.width = Pt(2)
        
        # Add text frame inside the content shape
        text_frame = content_shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        
        # Use full content if available, otherwise fall back to subheaders
        content_items = slide_data.get('content', slide_data.get('subheaders', []))
        
        # Log what we're working with for debugging
        logger.info(f"Slide: {slide_data.get('header', 'Unknown')}")
        logger.info(f"Has content: {bool(slide_data.get('content'))}")
        logger.info(f"Has subheaders: {bool(slide_data.get('subheaders'))}")
        logger.info(f"Content items count: {len(content_items) if isinstance(content_items, list) else 0}")
        
        # If we have content array, use it; otherwise use subheaders
        if isinstance(content_items, list) and len(content_items) > 0:
            # Process ALL content items, not just the first one
            for i, content_item in enumerate(content_items):
                if content_item and content_item.strip():  # Skip empty items
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # Clean up bullet points
                    text = content_item.strip()
                    if not text.startswith('•'):
                        text = f"• {text}"
                    
                    p.text = text
                    p.font.name = fonts['font_name']
                    p.font.size = Pt(fonts['content_size'])
                    p.font.color.rgb = colors['text_color']
                    p.level = 0
                    p.space_after = Pt(12)  # Add spacing between items
                    
                    logger.info(f"Added content item {i}: {text[:50]}...")
        else:
            # Fallback to subheaders if no content
            subheaders = slide_data.get('subheaders', [])
            logger.info(f"Using subheaders fallback: {len(subheaders)} items")
            
            for i, item in enumerate(subheaders):
                if item and item.strip():  # Skip empty items
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # Clean up bullet points
                    text = item.strip()
                    if not text.startswith('•'):
                        text = f"• {text}"
                    
                    p.text = text
                    p.font.name = fonts['font_name']
                    p.font.size = Pt(fonts['content_size'])
                    p.font.color.rgb = colors['text_color']
                    p.level = 0
                    p.space_after = Pt(12)  # Add spacing between items
                    
                    logger.info(f"Added subheader item {i}: {text[:50]}...")
        
        # Add decorative elements
        self.add_content_decorations(slide, colors)

    def save_presentation(self, prs, user_id, ai_result_id):
        """Save presentation to storage directory"""
        try:
            # Create storage directory if it doesn't exist
            storage_dir = os.path.join(os.path.dirname(__file__), '..', 'storage', 'app', 'presentations')
            os.makedirs(storage_dir, exist_ok=True)
            
            # Generate filename
            filename = f"presentation_{user_id}_{ai_result_id}_{int(time.time())}.pptx"
            filepath = os.path.join(storage_dir, filename)
        
            # Save presentation
            prs.save(filepath)
            
            logger.info(f"Presentation saved to: {filepath}")
            return filepath
        
        except Exception as e:
            logger.error(f"Failed to save presentation: {str(e)}")
            raise

def main():
    """Main entry point"""
    if len(sys.argv) != 2:
        print(json.dumps({
            'success': False,
            'error': 'Usage: python generate_presentation.py <data_file>'
        }))
        sys.exit(1)
    
    data_file = sys.argv[1]
    
    try:
        # Read data from file (handle UTF-8 BOM)
        with open(data_file, 'r', encoding='utf-8-sig') as f:
            data = json.load(f)
        
        # Generate presentation
        generator = PresentationGenerator()
        result = generator.generate_presentation(data)
        
        # Output result as JSON
        print(json.dumps(result))
        
    except Exception as e:
        logger.error(f"Script execution failed: {str(e)}")
        print(json.dumps({
            'success': False,
            'error': str(e)
        }))
        sys.exit(1)

if __name__ == '__main__':
    main()