#!/usr/bin/env python3
# Python PowerPoint Text Extractor Service
# Extracts text from PowerPoint presentations (.ppt, .pptx) using python-pptx
import sys
import json
import os

try:
    from pptx import Presentation
    import pptx
except ImportError as e:
    # When imported as module, don't exit, just raise the error
    if __name__ == "__main__":
        print(json.dumps({"success": False, "error": f"Missing required Python packages: {str(e)}", "text": "", "slides": 0, "metadata": {}, "word_count": 0, "character_count": 0}))
        sys.exit(1)
    else:
        raise ImportError(f"Missing required Python packages: {str(e)}")

def extract_text_from_slides(presentation):
    """Extract text from all slides"""
    text = ""
    slide_texts = []
    
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_text = ""
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text += shape.text + "\n"
            elif hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text.strip():
                        slide_text += paragraph.text + "\n"
        
        if slide_text.strip():
            slide_texts.append({
                "slide_number": slide_num,
                "text": slide_text.strip()
            })
            text += f"--- Slide {slide_num} ---\n{slide_text}\n\n"
    
    return text, slide_texts

def extract_text_from_notes(presentation):
    """Extract speaker notes from slides"""
    notes_text = ""
    notes_list = []
    
    for slide_num, slide in enumerate(presentation.slides, 1):
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                note_text = notes_slide.notes_text_frame.text.strip()
                if note_text:
                    notes_list.append({
                        "slide_number": slide_num,
                        "notes": note_text
                    })
                    notes_text += f"--- Notes for Slide {slide_num} ---\n{note_text}\n\n"
    
    return notes_text, notes_list

def extract_metadata_from_ppt(presentation):
    """Extract metadata from PowerPoint presentation"""
    try:
        core_props = presentation.core_properties
        metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "keywords": core_props.keywords or "",
            "comments": core_props.comments or "",
            "created": str(core_props.created) if core_props.created else "",
            "modified": str(core_props.modified) if core_props.modified else "",
            "last_modified_by": core_props.last_modified_by or "",
            "revision": core_props.revision or 0,
            "slide_count": len(presentation.slides)
        }
        return metadata
    except Exception as e:
        return {"error": str(e), "slide_count": len(presentation.slides)}

def extract_text_from_ppt(filepath):
    """Extract text from PowerPoint presentation"""
    try:
        presentation = Presentation(filepath)
        
        # Extract text from slides
        slides_text, slide_texts = extract_text_from_slides(presentation)
        
        # Extract speaker notes
        notes_text, notes_list = extract_text_from_notes(presentation)
        
        # Combine all text
        full_text = slides_text + notes_text
        
        # Get metadata
        metadata = extract_metadata_from_ppt(presentation)
        
        # Count words and characters
        word_count = len(full_text.split())
        character_count = len(full_text)
        
        return {
            "success": True,
            "text": full_text,
            "slides": len(presentation.slides),
            "slide_texts": slide_texts,
            "notes": notes_list,
            "metadata": metadata,
            "word_count": word_count,
            "character_count": character_count,
            "extraction_method": "python-pptx"
        }
        
    except Exception as e:
        return {
            "success": False,
            "text": "",
            "slides": 0,
            "metadata": {},
            "word_count": 0,
            "character_count": 0,
            "error": f"PowerPoint extraction failed: {str(e)}"
        }

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "error": "No file path provided", "text": "", "slides": 0, "metadata": {}, "word_count": 0, "character_count": 0}))
        sys.exit(1)
    
    ppt_filepath = sys.argv[1]
    if not os.path.exists(ppt_filepath):
        print(json.dumps({"success": False, "error": f"File not found: {ppt_filepath}", "text": "", "slides": 0, "metadata": {}, "word_count": 0, "character_count": 0}))
        sys.exit(1)
    
    result = extract_text_from_ppt(ppt_filepath)
    print(json.dumps(result, indent=2))
